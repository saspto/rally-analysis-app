"""
Rally Analysis App — Innovation Demo PowerPoint (12 slides, 16:9 widescreen).
Matches the style of the AWS Features Tracker deck.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ASSETS = "/workshop/rally-analysis-app/slides_assets"
OUT    = "/workshop/rally-analysis-app/Rally_Analysis_App_Innovation_Demo.pptx"

# ── Brand colours ─────────────────────────────────────────────────────────────
PRIMARY_700 = RGBColor(0x1E, 0x3A, 0x5F)   # header / buttons
PRIMARY_800 = RGBColor(0x16, 0x2D, 0x4A)   # active nav / dark bg
PRIMARY_600 = RGBColor(0x2D, 0x5F, 0x9E)   # accent
PRIMARY_100 = RGBColor(0xC5, 0xD5, 0xEA)   # light accent
PRIMARY_50  = RGBColor(0xE8, 0xEE, 0xF7)   # card tint

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF9, 0xFA, 0xFB)
MID_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY   = RGBColor(0x1E, 0x29, 0x3B)
BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_SEC    = RGBColor(0x64, 0x74, 0x8B)
SLATE_L     = RGBColor(0xF1, 0xF5, 0xF9)

GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_D     = RGBColor(0x14, 0x53, 0x2D)
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)
RED         = RGBColor(0xEF, 0x44, 0x44)
RED_D       = RGBColor(0x99, 0x1B, 0x1B)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
YELLOW      = RGBColor(0xCA, 0x8A, 0x04)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
TEAL        = RGBColor(0x0F, 0xB5, 0xD4)

# ── Slide dimensions — standard widescreen 16:9 ───────────────────────────────
SW = Emu(12192000)
SH = Emu(6858000)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def I(v): return Emu(int(v * 914400))
def P(v): return Pt(v)

def rect(slide, l, t, w, h, fill=None, line=None, line_pt=0.75):
    sh = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = P(line_pt)
    else:
        sh.line.fill.background()
    return sh

def txb(slide, l, t, w, h, text, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = P(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet(slide, l, t, w, h, text, size=11, dot_color=PRIMARY_700, text_color=DARK_GRAY):
    dot = I(0.1)
    rect(slide, l, t + I(size * 0.018), dot, dot, fill=dot_color)
    txb(slide, l + I(0.18), t, w - I(0.18), h, text, size, color=text_color)

def header_bar(slide, bg=PRIMARY_800, height=I(1.0)):
    rect(slide, 0, 0, SW, height, fill=bg)

def bottom_accent(slide, color=PRIMARY_600, height=I(0.1)):
    rect(slide, 0, SH - height, SW, height, fill=color)

def section_tag(slide, text, l=I(0.55), t=I(0.22)):
    txb(slide, l, t, I(6), I(0.28), text, 9, bold=True, color=PRIMARY_100)

def card(slide, l, t, w, h, fill=WHITE, border=BORDER):
    rect(slide, l + I(0.03), t + I(0.03), w, h, fill=RGBColor(0xCB, 0xD5, 0xE1))
    rect(slide, l, t, w, h, fill=fill, line=border, line_pt=0.75)

def screenshot_slide(title, section, img_path, points):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, SH, fill=PRIMARY_800)
    bottom_accent(sl)
    section_tag(sl, section)
    txb(sl, I(0.35), I(0.3), I(9.5), I(0.52), title, 20, bold=True, color=WHITE)

    sl.shapes.add_picture(img_path,
                           int(I(0.35)), int(I(0.92)),
                           int(I(8.38)), int(I(6.46)))

    PNL_L = I(8.85)
    PNL_T = I(0.92)
    PNL_W = I(3.12)
    PNL_H = I(6.46)
    rect(sl, PNL_L, PNL_T, PNL_W, PNL_H, fill=RGBColor(0x0A, 0x14, 0x24))
    rect(sl, PNL_L, PNL_T, I(0.06), PNL_H, fill=PRIMARY_600)

    txb(sl, PNL_L + I(0.18), PNL_T + I(0.18), PNL_W - I(0.28), I(0.28),
        "KEY HIGHLIGHTS", 8, bold=True, color=PRIMARY_100)

    for j, pt in enumerate(points):
        py = PNL_T + I(0.58 + j * 0.94)
        rect(sl, PNL_L + I(0.18), py, PNL_W - I(0.28), I(0.82),
             fill=RGBColor(0x16, 0x28, 0x40))
        rect(sl, PNL_L + I(0.18), py, I(0.04), I(0.82), fill=PRIMARY_600)
        txb(sl, PNL_L + I(0.28), py + I(0.1), PNL_W - I(0.42), I(0.65),
            pt, 10, color=RGBColor(0xE2, 0xE8, 0xF0))
    return sl


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Hero
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=PRIMARY_800)
rect(sl, 0, 0, I(0.18), SH, fill=PRIMARY_600)
for i in range(1, 9):
    rect(sl, I(i * 13.333/8), 0, I(0.02), SH, fill=RGBColor(0x1A, 0x30, 0x4E))

rect(sl, I(0.65), I(1.05), I(0.72), I(0.48), fill=PRIMARY_600)
txb(sl, I(0.67), I(1.08), I(0.68), I(0.42), "RA", 18, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

txb(sl, I(0.65), I(1.7), I(10), I(0.3),
    "INNOVATION DEMO  ·  ENGINEERING & PROJECT MANAGEMENT",
    10, bold=True, color=PRIMARY_100)

txb(sl, I(0.65), I(2.3), I(11), I(1.0),
    "Rally Analysis App", 52, bold=True, color=WHITE)

txb(sl, I(0.65), I(3.45), I(10), I(0.5),
    "AI-Powered Sprint Analytics & Executive Reporting for Rally / CA Agile Central",
    21, italic=True, color=RGBColor(0xCB, 0xD5, 0xE1))

for i, lbl in enumerate(["Automated daily exports",
                          "Claude AI summaries",
                          "Sprint analytics & insights"]):
    px = I(0.65 + i * 3.3)
    rect(sl, px, I(4.25), I(3.05), I(0.42), fill=PRIMARY_600)
    txb(sl, px, I(4.27), I(3.05), I(0.38), lbl, 12,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 0, SH - I(0.48), SW, I(0.48), fill=RGBColor(0x0A, 0x18, 0x2E))
txb(sl, 0, SH - I(0.44), SW, I(0.4),
    "Presented by the Engineering & Cloud Administration Team  ·  April 2026",
    10, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=LIGHT_GRAY)
header_bar(sl, height=I(1.05))
bottom_accent(sl)

txb(sl, I(0.55), I(0.18), I(10.5), I(0.56),
    "The Challenge: Rally Data Without Context", 28, bold=True, color=WHITE)
txb(sl, I(0.55), I(0.74), I(10.5), I(0.28),
    "Rally holds all your project data — but generating insights requires hours of manual work.",
    13, italic=True, color=PRIMARY_100)

for i, (num, lbl) in enumerate([("Hours",  "spent weekly writing\nsprint summaries"),
                                  ("3+",     "tools open to answer\none status question"),
                                  ("0",      "minutes execs spend\nreading raw CSV data"),
                                  ("Manual", "copy-paste process\nfrom Rally to reports")]):
    bx = I(0.55 + i * 3.12)
    card(sl, bx, I(1.18), I(2.9), I(1.52))
    txb(sl, bx, I(1.28), I(2.9), I(0.72), num, 36, bold=True,
        color=PRIMARY_700, align=PP_ALIGN.CENTER)
    txb(sl, bx, I(2.02), I(2.9), I(0.58), lbl, 12,
        color=DARK_GRAY, align=PP_ALIGN.CENTER)

rect(sl, I(0.55), I(2.88), SW - I(1.1), I(0.01), fill=BORDER)
txb(sl, I(0.55), I(2.97), I(5.5), I(0.32), "Pain Points", 15, bold=True, color=DARK_GRAY)

for i, pain in enumerate([
    "Generating sprint summaries is a 2–3 hour manual process at end of every sprint",
    "Rally exports are raw CSV — no narrative context, no insights, no executive language",
    "Analytics require pivot tables, formulas, and manual chart building in Excel",
    "Stale or blocked stories go unnoticed until stand-ups — too late to act",
    "Different stakeholders need different formats: engineering vs. C-suite vs. PM",
]):
    bullet(sl, I(0.55), I(3.42 + i * 0.5), I(8.2), I(0.46), pain, 12,
           dot_color=RED, text_color=DARK_GRAY)

card(sl, I(9.15), I(2.97), I(3.62), I(3.42))
txb(sl, I(9.28), I(3.08), I(3.36), I(0.34), "The Opportunity",
    14, bold=True, color=PRIMARY_700)
txb(sl, I(9.28), I(3.48), I(3.36), I(2.75),
    "Automate every step — daily export, AI narrative generation, analytics computation "
    "— so your team gets insights without the manual grind.",
    12, color=DARK_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=PRIMARY_800)
bottom_accent(sl)
section_tag(sl, "SOLUTION ARCHITECTURE")
txb(sl, I(0.55), I(0.48), I(10), I(0.52),
    "Rally Analysis App — How It Works", 27, bold=True, color=WHITE)

FLOW_Y = I(1.32)
BOX_W  = I(2.18)
BOX_H  = I(1.1)

ingestion = [
    ("Rally\n(CA Agile Central)", "Features, Stories,\nTasks via API",   PRIMARY_600),
    ("Fetcher\nLambda",           "Daily 3PM EST\n+ on-demand",          BLUE),
    ("S3 Exports\nBucket",        "CSV per day\nwith presigned URLs",    TEAL),
    ("Summary\nLambda",           "Claude Opus 4\n3 summary types",      GREEN),
    ("S3 Summaries\n+ API",       "JSON summaries\n+ analytics compute", PURPLE),
]
for i, (name, desc, color) in enumerate(ingestion):
    bx = I(0.38 + i * (2.18 + 0.28))
    rect(sl, bx, FLOW_Y, BOX_W, BOX_H, fill=color)
    txb(sl, bx, FLOW_Y + I(0.1), BOX_W, I(0.44), name, 12,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, bx, FLOW_Y + I(0.55), BOX_W, I(0.48), desc, 10,
        color=RGBColor(0xE2, 0xE8, 0xF0), align=PP_ALIGN.CENTER)
    if i < len(ingestion) - 1:
        ax = bx + BOX_W + I(0.06)
        txb(sl, ax, FLOW_Y + I(0.38), I(0.18), I(0.3), "→", 18,
            bold=True, color=PRIMARY_100, align=PP_ALIGN.CENTER)

txb(sl, I(0.38), I(2.65), I(12), I(0.28),
    "▼  Presentation Layer", 10, bold=True, color=MID_GRAY)

API_Y = I(3.05)
api = [
    ("EventBridge\nScheduler",  "Daily cron\n3 PM EST",              ORANGE),
    ("API Gateway\nHTTP API",   "5 endpoints\nCORS enabled",         BLUE),
    ("CloudFront CDN",          "Global delivery\n+ OAC for S3",     PRIMARY_600),
    ("React SPA (S3)",          "TypeScript + Vite\nTailwind CSS",   PURPLE),
]
for i, (name, desc, color) in enumerate(api):
    bx = I(0.38 + i * (2.18 + 0.28))
    rect(sl, bx, API_Y, BOX_W, BOX_H, fill=color)
    txb(sl, bx, API_Y + I(0.1), BOX_W, I(0.44), name, 12,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, bx, API_Y + I(0.55), BOX_W, I(0.48), desc, 10,
        color=RGBColor(0xE2, 0xE8, 0xF0), align=PP_ALIGN.CENTER)
    if i < len(api) - 1:
        ax = bx + BOX_W + I(0.06)
        txb(sl, ax, API_Y + I(0.38), I(0.18), I(0.3), "→", 18,
            bold=True, color=PRIMARY_100, align=PP_ALIGN.CENTER)

txb(sl, I(0.38), I(4.38), I(12), I(0.26), "Tech Stack", 10, bold=True, color=MID_GRAY)
for i, t in enumerate(["Python 3.12 · Lambda arm64",
                        "Anthropic Claude Opus 4",
                        "Rally REST API v2",
                        "React 18 + TypeScript + Vite",
                        "Recharts · jsPDF · docx · xlsx",
                        "AWS SAM (Infrastructure as Code)"]):
    bx = I(0.38 + i * 2.14)
    rect(sl, bx, I(4.7), I(2.04), I(0.4), fill=RGBColor(0x16, 0x2D, 0x4A))
    txb(sl, bx, I(4.72), I(2.04), I(0.36), t, 9,
        color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Features
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=LIGHT_GRAY)
header_bar(sl, height=I(1.05))
bottom_accent(sl)

txb(sl, I(0.55), I(0.18), I(11), I(0.56),
    "Everything Your Team Needs — Automated", 27, bold=True, color=WHITE)
txb(sl, I(0.55), I(0.74), I(11), I(0.28),
    "From daily Rally exports to AI-generated executive reports — zero manual effort.",
    13, italic=True, color=PRIMARY_100)

features = [
    (PRIMARY_600, "Automated Daily Exports",
     "EventBridge triggers a Lambda every day at 3 PM EST to pull Features, Stories, and Tasks from Rally. CSV saved to S3 with presigned download URLs.",
     ["Zero manual export steps ever again", "On-demand trigger also available in UI"]),
    (GREEN, "Three AI Summary Types",
     "Claude Opus 4 generates Weekly, Monthly, and Quarterly summaries — each tailored to audience. Main narrative, executive 2-pager, and per-story detail.",
     ["Main summary: 4,096 tokens of narrative", "Executive: C-suite focused, metrics-first"]),
    (BLUE, "Sprint Analytics Dashboard",
     "Automatic computation of burndown, velocity, cycle time, and resource allocation from raw Rally data. No spreadsheets needed.",
     ["Burn charts: actual vs. ideal per sprint", "Velocity trend across rolling sprints"]),
    (ORANGE, "Stale Story Detection",
     "Automatically surfaces stories not updated in 7+ days. Warning at 7d, critical at 14d. Sorted by severity so the most blocked items appear first.",
     ["Warning threshold: 7 days without update", "Critical threshold: 14+ days — acts as blocker"]),
    (PURPLE, "Multi-Format Export",
     "Download any summary as PDF, Word DOCX, or Excel XLSX with a single click. Each format is properly structured with headings, bullets, and tables.",
     ["PDF via jsPDF with markdown rendering", "DOCX via docx library — ready to email"]),
    (YELLOW, "Cycle Time & Outlier Detection",
     "Tracks story creation to acceptance date. Flags statistical outliers at mean + 1.5σ so teams can investigate systemic delays proactively.",
     ["Mean + 1.5σ outlier threshold", "Owner-level breakdown for accountability"]),
]

COL_W = I(4.05)
ROW_H = I(1.82)
for i, (color, title, body, pts) in enumerate(features):
    col = i % 3
    row = i // 3
    bx = I(0.42 + col * (4.05 + 0.12))
    by = I(1.28 + row * (1.82 + 0.1))
    card(sl, bx, by, COL_W, ROW_H)
    rect(sl, bx, by, I(0.08), ROW_H, fill=color)
    txb(sl, bx + I(0.15), by + I(0.1), COL_W - I(0.22), I(0.3),
        title, 13, bold=True, color=DARK_GRAY)
    txb(sl, bx + I(0.15), by + I(0.44), COL_W - I(0.22), I(0.52),
        body, 9.5, color=TEXT_SEC)
    for j, pt in enumerate(pts[:2]):
        bullet(sl, bx + I(0.15), by + I(1.0 + j * 0.32), COL_W - I(0.22), I(0.28),
               pt, 9, dot_color=color, text_color=DARK_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDES 5–7 — Screenshots
# ═════════════════════════════════════════════════════════════════════════════
screenshot_slide(
    title="Exports Dashboard — Daily CSV Exports from Rally",
    section="LIVE DEMO — EXPORTS DASHBOARD",
    img_path=f"{ASSETS}/screen_exports.png",
    points=[
        "Auto-export runs\ndaily at 3 PM EST",
        "Manual trigger\navailable on demand",
        "Feature / Story / Task\ncounts per export",
        "Presigned S3 URLs\nfor secure download",
        "Full export history\nretained in S3",
        "CSV format for\ndownstream tooling",
    ],
)

screenshot_slide(
    title="Analytics Dashboard — Sprint Burndown, Velocity & Stale Stories",
    section="LIVE DEMO — ANALYTICS DASHBOARD",
    img_path=f"{ASSETS}/screen_analytics.png",
    points=[
        "6 KPI cards: stories,\npoints, velocity, cycle time",
        "Burndown: actual vs\nideal line per sprint",
        "Three sprints shown:\nApr / May / Jun 2026",
        "Resource allocation\nby team member",
        "Stale story table\nwith severity flags",
        "Outlier detection\nat mean + 1.5 sigma",
    ],
)

screenshot_slide(
    title="AI Summary Page — Claude-Generated Narrative Summaries",
    section="LIVE DEMO — AI SUMMARY GENERATION",
    img_path=f"{ASSETS}/screen_summary.png",
    points=[
        "Date range picker +\nWeekly/Monthly/Quarterly",
        "6 live metrics from\nRally export data",
        "Full narrative summary\ngenerated by Claude",
        "One-click export:\nPDF, DOCX, or Excel",
        "Navigate to Executive\nor Detailed summary",
        "Structured markdown:\nOverview, Risks, Metrics",
    ],
)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Use Cases
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=LIGHT_GRAY)
header_bar(sl, height=I(1.05))
bottom_accent(sl)

txb(sl, I(0.55), I(0.18), I(11), I(0.56),
    "Who Uses It and How", 27, bold=True, color=WHITE)
txb(sl, I(0.55), I(0.74), I(11), I(0.28),
    "Every role in the delivery chain gets what they need — automatically.",
    13, italic=True, color=PRIMARY_100)

use_cases = [
    (PRIMARY_600, "Engineering Manager",
     "Every Friday afternoon, open the app, select the last 2-week sprint, and generate a "
     "Monthly summary. Download as DOCX and paste into the sprint retrospective doc in 2 minutes.",
     "Summary type: Monthly · Download: DOCX"),
    (GREEN, "C-Suite / Executive",
     "Receive a weekly Executive Summary email (Slack/SNS roadmap item) with portfolio health "
     "status, velocity trend, and any active blockers — no Rally login required.",
     "Executive Summary · Metrics · Health status"),
    (ORANGE, "Scrum Master / Agile Coach",
     "Open Analytics before daily stand-up to check the stale stories table. Immediately see "
     "which stories have gone silent and who owns them.",
     "Analytics → Stale & Blocked Stories"),
    (BLUE, "Portfolio Manager",
     "Select a full quarter date range and generate a Quarterly summary. Download as PDF for "
     "inclusion in the QBR deck — done in under 5 minutes.",
     "Summary type: Quarterly · Download: PDF"),
    (PURPLE, "Release Manager",
     "Use the cycle time chart to identify which story types consistently become outliers. "
     "Feed findings into the Definition of Ready to prevent future delays.",
     "Analytics → Cycle Time Analysis"),
    (RED, "Engineering Lead",
     "Review the Resource Allocation chart to spot team members who are overloaded (high "
     "actual vs. estimated hours) or under-utilized before the next sprint planning.",
     "Analytics → Resource Allocation"),
]

COL_W = I(4.05)
ROW_H = I(2.02)
for i, (color, title, body, query) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    bx = I(0.42 + col * (4.05 + 0.12))
    by = I(1.28 + row * (2.02 + 0.06))
    card(sl, bx, by, COL_W, ROW_H)
    rect(sl, bx, by, COL_W, I(0.07), fill=color)
    txb(sl, bx + I(0.12), by + I(0.13), COL_W - I(0.2), I(0.3),
        title, 13, bold=True, color=DARK_GRAY)
    txb(sl, bx + I(0.12), by + I(0.46), COL_W - I(0.2), I(0.88),
        body, 9.5, color=TEXT_SEC)
    rect(sl, bx + I(0.12), by + ROW_H - I(0.38), COL_W - I(0.24), I(0.28),
         fill=SLATE_L)
    txb(sl, bx + I(0.18), by + ROW_H - I(0.36), COL_W - I(0.3), I(0.26),
        query, 8.5, italic=True, color=color)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI Deep-Dive
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=PRIMARY_800)
bottom_accent(sl)
section_tag(sl, "AI INTELLIGENCE LAYER")
txb(sl, I(0.55), I(0.42), I(10.5), I(0.52),
    "How Claude Opus 4 Transforms Raw Rally Data Into Insights",
    25, bold=True, color=WHITE)

rect(sl, I(0.45), I(1.08), I(5.5), I(5.88),
     fill=RGBColor(0x0A, 0x18, 0x2E), line=RGBColor(0x1E, 0x3A, 0x5F))
txb(sl, I(0.6), I(1.14), I(5.2), I(0.28), "RAW RALLY EXPORT (CSV)", 9,
    bold=True, color=MID_GRAY)
txb(sl, I(0.6), I(1.48), I(5.2), I(5.3),
    "FormattedID,Name,State,Owner,Points,CreationDate,AcceptedDate\n"
    "US-001,User auth flow,Accepted,Alice Johnson,8,2026-04-01,2026-04-09\n"
    "US-002,Dashboard layout,Accepted,Bob Smith,5,2026-04-01,2026-04-08\n"
    "US-003,Data export CSV,Accepted,Alice Johnson,13,2026-04-02,2026-04-10\n"
    "US-004,OAuth integration,In-Progress,Carol Davis,5,2026-04-03,\n"
    "US-005,Search/filter UI,Completed,Bob Smith,3,2026-04-04,\n"
    "US-006,Real-time refresh,In-Progress,Eve Martinez,8,2026-03-01,\n"
    "US-007,API rate limiting,In-Progress,Grace Kim,5,2026-03-15,\n"
    "US-008,Webhook events,Defined,Henry Wang,8,2026-02-28,\n\n"
    "F-001,Customer Portal,Accepted,,,2026-03-15,2026-04-10\n"
    "F-002,API Integration,In-Progress,,,2026-03-20,\n"
    "F-003,Performance Opt.,Defined,,,2026-04-01,",
    8.5, color=RGBColor(0x7D, 0xD3, 0xFC))

txb(sl, I(6.12), I(3.55), I(0.9), I(0.5), "→", 36,
    bold=True, color=PRIMARY_600, align=PP_ALIGN.CENTER)
txb(sl, I(5.98), I(4.12), I(1.18), I(0.36), "Claude\nOpus 4", 9,
    color=PRIMARY_100, align=PP_ALIGN.CENTER)

rect(sl, I(7.22), I(1.08), I(5.5), I(5.88),
     fill=RGBColor(0x0A, 0x20, 0x14), line=GREEN)
txb(sl, I(7.38), I(1.14), I(5.2), I(0.28), "AI-GENERATED SUMMARY OUTPUT", 9,
    bold=True, color=MID_GRAY)

ty = I(1.48)
for label, val, col in [
    ("summary type:", "weekly", PRIMARY_100),
    ("narrative:", (
        "The team delivered strong results this week,\n"
        "completing 5 user stories (34 story points)\n"
        "across the Customer Portal and API Integration\n"
        "features. Feature F-001 reached full acceptance.\n"
        "Team velocity of 38 pts exceeds the 30-pt target."
    ), RGBColor(0xBB, 0xF7, 0xD0)),
    ("metrics:", (
        "features_completed: 1\n"
        "stories_completed:  5\n"
        "stories_in_progress: 3\n"
        "tasks_completed:    11\n"
        "total_points:       34\n"
        "team_velocity:      38"
    ), RGBColor(0x86, 0xEF, 0xAC)),
    ("risks:", (
        "US-008 (Henry Wang) — 45 days stale\n"
        "Recommend escalation before next sprint."
    ), RGBColor(0xFD, 0xBA, 0x74)),
]:
    txb(sl, I(7.38), ty, I(5.2), I(0.26), label, 9, bold=True, color=MID_GRAY)
    ty += I(0.28)
    txb(sl, I(7.38), ty, I(5.2), I(1.0), val, 10, color=col)
    ty += I(0.88)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Cost & Security
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=LIGHT_GRAY)
header_bar(sl, height=I(1.05))
bottom_accent(sl)

txb(sl, I(0.55), I(0.18), I(11), I(0.56),
    "Lightweight, Secure, and Production-Ready", 27, bold=True, color=WHITE)
txb(sl, I(0.55), I(0.74), I(11), I(0.28),
    "Serverless architecture keeps costs minimal. IAM least-privilege, S3 OAC, no public data exposure.",
    13, italic=True, color=PRIMARY_100)

card(sl, I(0.45), I(1.18), I(5.82), I(4.88))
txb(sl, I(0.6), I(1.28), I(5.5), I(0.3),
    "Monthly Cost Estimate (~500 summaries/month)", 12, bold=True, color=DARK_GRAY)

rows = [
    ("Service", "Monthly Cost", True),
    ("Lambda (4 functions, 512MB, 300s)", "~$0.05", False),
    ("API Gateway HTTP API", "~$0.01", False),
    ("S3 (exports + summaries + SPA)", "~$0.05", False),
    ("CloudFront (PriceClass_100)", "~$0.10", False),
    ("Anthropic Claude Opus 4 (500 calls)", "~$7.50", False),
    ("EventBridge Scheduler", "$0.00  (free tier)", False),
    ("TOTAL", "~$7.71/month", True),
]
for i, (svc, cost, bold) in enumerate(rows):
    ry = I(1.7 + i * 0.46)
    bg = RGBColor(0xF8, 0xFA, 0xFC) if (i % 2 == 0 and not bold) else None
    if bold:
        bg = PRIMARY_50
    if bg:
        rect(sl, I(0.45), ry, I(5.82), I(0.46), fill=bg)
    txb(sl, I(0.6), ry + I(0.1), I(3.8), I(0.28), svc, 10,
        bold=bold, color=DARK_GRAY)
    txb(sl, I(4.55), ry + I(0.1), I(1.5), I(0.28), cost, 10,
        bold=bold, color=PRIMARY_700 if bold else GREEN, align=PP_ALIGN.RIGHT)

card(sl, I(6.62), I(1.18), I(6.12), I(2.25))
txb(sl, I(6.78), I(1.28), I(5.8), I(0.3),
    "Security & Compliance", 13, bold=True, color=DARK_GRAY)
for j, item in enumerate([
    "S3 Block Public Access ON — CloudFront OAC for frontend, no direct bucket access",
    "IAM least-privilege — each Lambda has only the permissions it needs",
    "Rally API key and Anthropic key stored in SSM Parameter Store (encrypted)",
    "API Gateway no-auth read endpoints — no write surface exposed publicly",
    "All traffic via HTTPS — CloudFront enforces TLS, no HTTP",
]):
    bullet(sl, I(6.78), I(1.68 + j * 0.3), I(5.8), I(0.28),
           item, 9.5, dot_color=GREEN, text_color=DARK_GRAY)

card(sl, I(6.62), I(3.55), I(6.12), I(2.52))
txb(sl, I(6.78), I(3.65), I(5.8), I(0.3),
    "Reliability & Observability", 13, bold=True, color=DARK_GRAY)
for j, item in enumerate([
    "EventBridge schedule survives Lambda cold-starts — retries built in",
    "Lambda timeout 300s — handles large Rally workspaces without truncation",
    "Mock mode (VITE_USE_MOCK=true) enables full UI demo without AWS/Rally",
    "Presigned S3 URLs expire after 1 hour — no persistent public links",
    "Python pytest + Vitest unit test suites for backend and frontend",
]):
    bullet(sl, I(6.78), I(4.08 + j * 0.3), I(5.8), I(0.28),
           item, 9.5, dot_color=BLUE, text_color=DARK_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Roadmap
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=PRIMARY_800)
bottom_accent(sl)
section_tag(sl, "WHAT'S NEXT")
txb(sl, I(0.55), I(0.42), I(10), I(0.52),
    "Roadmap — Where We Can Take This", 25, bold=True, color=WHITE)

roadmap = [
    ("Q2 2026", PRIMARY_600, [
        "Slack / Teams integration — post summary to #sprint-updates automatically after each export",
        "Multi-project support — select from multiple Rally workspaces in a single dashboard",
        "Trend charts — rolling 6-month velocity and cycle time trend visualization",
    ]),
    ("Q3 2026", GREEN, [
        "Predictive risk scoring — Claude flags stories at risk of missing sprint before it happens",
        "Jira / Azure DevOps connector — same analytics for teams not on Rally",
        "Shareable report links — generate a read-only URL for any summary to share with stakeholders",
    ]),
    ("Q4 2026", ORANGE, [
        "Executive portfolio view — multi-team rollup with cross-project health scores",
        "Capacity planning assistant — Claude recommends sprint load based on historical velocity",
        "Auto-generated retrospective — AI identifies patterns and suggests process improvements",
    ]),
]

for i, (quarter, color, items) in enumerate(roadmap):
    bx = I(0.42 + i * 4.28)
    by = I(1.05)
    bw = I(4.05)
    bh = I(5.68)
    rect(sl, bx, by, bw, bh,
         fill=RGBColor(0x0E, 0x1E, 0x32), line=color, line_pt=1.5)
    rect(sl, bx, by, bw, I(0.44), fill=color)
    txb(sl, bx, by + I(0.06), bw, I(0.34), quarter, 16,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        iy = by + I(0.58 + j * 1.62)
        rect(sl, bx + I(0.22), iy, I(3.62), I(1.46),
             fill=RGBColor(0x16, 0x2A, 0x3E))
        bullet(sl, bx + I(0.34), iy + I(0.12), I(3.38), I(1.2),
               item, 11, dot_color=color,
               text_color=RGBColor(0xCB, 0xD5, 0xE1))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Call to Action
# ═════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, fill=PRIMARY_800)
rect(sl, 0, 0, I(0.18), SH, fill=PRIMARY_600)
for i in range(1, 9):
    rect(sl, I(i * 13.333/8), 0, I(0.02), SH, fill=RGBColor(0x1A, 0x30, 0x4E))

txb(sl, I(0.65), I(1.05), I(11.5), I(0.4),
    "The app is live and running today.", 23, italic=True, color=MID_GRAY)
txb(sl, I(0.65), I(1.6), I(11.5), I(1.0),
    "Your next sprint report writes itself.", 48, bold=True, color=WHITE)
txb(sl, I(0.65), I(2.75), I(11), I(0.42),
    "https://d28e9pzfl4v2ga.cloudfront.net",
    20, bold=True, color=PRIMARY_100)

for i, (color, title, body) in enumerate([
    (PRIMARY_600, "Try It Now",
     "Open the app, trigger a manual export from your Rally workspace, and generate your first AI summary in under 2 minutes."),
    (GREEN, "Connect Your Workspace",
     "Provide your Rally ZSESSIONID and workspace name — we'll configure the Lambda environment variables and run the first export for you."),
    (ORANGE, "Request a Feature",
     "Tell us which report formats, integrations, or analytics views would make this tool indispensable to your sprint workflow."),
]):
    bx = I(0.65 + i * 4.18)
    by = I(3.72)
    rect(sl, bx, by, I(3.9), I(2.48), fill=RGBColor(0x0E, 0x1E, 0x32))
    rect(sl, bx, by, I(3.9), I(0.07), fill=color)
    txb(sl, bx + I(0.18), by + I(0.18), I(3.6), I(0.32), title,
        14, bold=True, color=WHITE)
    txb(sl, bx + I(0.18), by + I(0.58), I(3.6), I(1.7), body,
        10.5, color=RGBColor(0xCB, 0xD5, 0xE1))

rect(sl, 0, SH - I(0.48), SW, I(0.48), fill=RGBColor(0x08, 0x12, 0x20))
txb(sl, 0, SH - I(0.44), SW, I(0.4),
    "Rally Analysis App  ·  Lambda · Claude Opus 4 · DynamoDB · CloudFront  ·  AWS SAM  ·  April 2026",
    9, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Fix sldSz type so PowerPoint accepts the file ─────────────────────────────
for el in prs._element.iter():
    tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
    if tag == 'sldSz':
        el.set('type', 'custom')
        break

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
